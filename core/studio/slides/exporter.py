"""PPTX renderer for Forge slides — programmatic shapes, no templates."""

import io
import re
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

import hashlib

from core.schemas.studio_schema import SlideTheme, SlidesContentTree
from core.studio.slides.themes import _blend_color

# 16:9 widescreen dimensions
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# Margins
MARGIN_LEFT = Inches(0.75)
MARGIN_TOP = Inches(0.75)
MARGIN_RIGHT = Inches(0.75)
MARGIN_BOTTOM = Inches(0.5)

# Content area
CONTENT_WIDTH = SLIDE_WIDTH - MARGIN_LEFT - MARGIN_RIGHT
CONTENT_HEIGHT = SLIDE_HEIGHT - MARGIN_TOP - MARGIN_BOTTOM

# Title area
TITLE_TOP = Inches(0.5)
TITLE_LEFT = MARGIN_LEFT
TITLE_WIDTH = CONTENT_WIDTH
TITLE_HEIGHT = Inches(1.0)

# Body area (below title)
BODY_TOP = Inches(1.8)
BODY_LEFT = MARGIN_LEFT
BODY_WIDTH = CONTENT_WIDTH
BODY_HEIGHT = Inches(5.0)

# Two-column split
COLUMN_GAP = Inches(0.5)
COLUMN_WIDTH = (CONTENT_WIDTH - COLUMN_GAP) / 2

# Chart area (used when slide_type == "chart")
CHART_TOP = Inches(2.0)
CHART_LEFT = MARGIN_LEFT
CHART_WIDTH = CONTENT_WIDTH
CHART_HEIGHT = Inches(3.8)

# Caption below chart
CAPTION_TOP = Inches(6.0)
CAPTION_HEIGHT = Inches(1.0)

# === Design Tokens ===

_DESIGN_TOKENS = {
    # Typography scale
    "title_size": Pt(48),
    "heading_size": Pt(36),
    "subheading_size": Pt(24),
    "body_size": Pt(18),
    "body_small_size": Pt(16),
    "code_size": Pt(14),
    "footer_size": Pt(10),
    "quote_size": Pt(32),
    "attribution_size": Pt(18),
    "stat_callout_size": Pt(64),
    "stat_label_size": Pt(16),

    # Text frame margins
    "margin_left": Inches(0.15),
    "margin_right": Inches(0.15),
    "margin_top": Inches(0.08),
    "margin_bottom": Inches(0.08),

    # Paragraph spacing
    "line_spacing": 1.15,
    "para_space_after": Pt(6),
    "bullet_space_after": Pt(8),

    # Bullet formatting
    "bullet_indent": Inches(0.25),
    "bullet_hanging": Inches(0.20),

    # Kicker / Takeaway
    "kicker_size": Pt(12),
    "takeaway_size": Pt(16),

    # Caption tokens
    "caption_size": Pt(14),

    # Slide chrome
    "accent_bar_height": Inches(0.10),
    "accent_bar_top": Inches(7.1),
    "footer_height": Inches(0.25),
    "footer_top": Inches(7.2),
}

# Font scale compensation for serif fonts
_FONT_SCALE = {
    "Garamond": 1.12,
    "Book Antiqua": 1.10,
    "Georgia": 1.08,
    "Constantia": 1.08,
}

# Chart type mapping
try:
    from pptx.chart.data import CategoryChartData, XyChartData
    from pptx.enum.chart import XL_CHART_TYPE
    _CHART_TYPE_MAP = {
        "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
        "line": XL_CHART_TYPE.LINE,
        "pie": XL_CHART_TYPE.PIE,
        "funnel": XL_CHART_TYPE.BAR_CLUSTERED,
    }
    _CHARTS_AVAILABLE = True
except ImportError:
    _CHARTS_AVAILABLE = False
    _CHART_TYPE_MAP = {}


def export_to_pptx(
    content_tree: SlidesContentTree,
    theme: SlideTheme,
    output_path: Path,
    images: dict[str, io.BytesIO] | None = None,
) -> Path:
    """Export a SlidesContentTree to PPTX format.

    Args:
        images: Optional dict mapping slide ID to JPEG BytesIO buffers for
                image_text slides. When provided, actual images are embedded
                instead of text placeholders.

    Returns the output file path.
    """
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    blank_layout = prs.slide_layouts[6]  # Blank layout
    total_slides = len(content_tree.slides)
    section_count = 0

    for i, slide_data in enumerate(content_tree.slides):
        pptx_slide = prs.slides.add_slide(blank_layout)

        if slide_data.slide_type == "section_divider":
            section_count += 1

        # Draw motif first (behind other shapes)
        _add_slide_motif(pptx_slide, theme)

        renderer = _RENDERERS.get(slide_data.slide_type, _render_content)
        renderer(pptx_slide, slide_data, theme, images=images, slide_index=i,
                 section_number=section_count)

        # Add slide chrome (skip first and last slides)
        if 0 < i < total_slides - 1:
            _add_slide_chrome(pptx_slide, theme, i + 1, total_slides)

        if slide_data.speaker_notes:
            notes_slide = pptx_slide.notes_slide
            notes_slide.notes_text_frame.text = slide_data.speaker_notes

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    return output_path


# === Helper Functions ===

def _scaled_font_size(base_size, font_name):
    """Apply font scale compensation for serif fonts."""
    scale = _FONT_SCALE.get(font_name, 1.0)
    if scale == 1.0:
        return base_size
    return Pt(int(base_size.pt * scale + 0.5))


def _compute_body_top(title: str, font_size_pt: float = 36.0) -> tuple:
    """Compute adjusted body top and title font size based on title length.

    Returns (body_top_inches, title_font_size_pt).
    Long titles that wrap to multiple lines push the body area down.
    """
    title = (title or "").strip()
    title_len = len(title)

    # Approximate characters per line at given font size on CONTENT_WIDTH (~11.83")
    # At 36pt ~2 chars/inch → ~24 chars/line; at 28pt ~2.5 chars/inch → ~30 chars/line
    if title_len > 80:
        # 3+ lines likely — shrink title font and push body down
        adjusted_font = 28.0
        body_top = 2.3
    elif title_len > 50:
        # 2 lines likely — push body down slightly
        adjusted_font = font_size_pt
        body_top = 2.1
    else:
        # Single line — default
        adjusted_font = font_size_pt
        body_top = 1.8

    # Ensure body area stays >= 4.2" (slide bottom at ~7.1" minus body_top)
    max_body_top = 7.1 - 4.2  # = 2.9"
    body_top = min(body_top, max_body_top)

    return Inches(body_top), Pt(int(adjusted_font))


_MD_INLINE_RE = re.compile(
    r'\*\*\*(?!\s)(.+?)(?<!\s)\*\*\*'   # ***bold+italic*** (no inner spaces at delimiters)
    r'|\*\*(?!\s)(.+?)(?<!\s)\*\*'       # **bold**
    r'|\*(?!\s)(.+?)(?<!\s)\*'           # *italic*
    r'|==(?!\s)(.+?)(?<!\s)=='           # ==highlight==
)


def _parse_markdown_runs(text: str) -> list[tuple[str, bool, bool, bool]]:
    """Parse markdown inline formatting into (text, bold, italic, highlight) segments."""
    if not text:
        return [("", False, False, False)]
    segments: list[tuple[str, bool, bool, bool]] = []
    last_end = 0
    for m in _MD_INLINE_RE.finditer(text):
        # Plain text before this match
        if m.start() > last_end:
            segments.append((text[last_end:m.start()], False, False, False))
        if m.group(1) is not None:       # ***bold+italic***
            segments.append((m.group(1), True, True, False))
        elif m.group(2) is not None:     # **bold**
            segments.append((m.group(2), True, False, False))
        elif m.group(3) is not None:     # *italic*
            segments.append((m.group(3), False, True, False))
        elif m.group(4) is not None:     # ==highlight==
            segments.append((m.group(4), False, False, True))
        last_end = m.end()
    # Trailing plain text (or entire string if no matches)
    if last_end < len(text):
        segments.append((text[last_end:], False, False, False))
    return segments or [("", False, False, False)]


def _apply_markdown_runs(paragraph, text, *, font_name, font_size, font_color,
                         bold=False, accent_color=None):
    """Replace paragraph text with per-run markdown formatting."""
    if isinstance(font_color, str):
        font_color = RGBColor.from_string(font_color.lstrip("#"))
    if isinstance(accent_color, str):
        accent_color = RGBColor.from_string(accent_color.lstrip("#"))

    segments = _parse_markdown_runs(str(text))

    # Clear existing text
    paragraph.clear()

    # Set paragraph-level font size for validator compatibility
    paragraph.font.size = font_size

    for seg_text, seg_bold, seg_italic, seg_highlight in segments:
        run = paragraph.add_run()
        run.text = seg_text
        run.font.name = font_name
        run.font.size = font_size
        if seg_highlight and accent_color is not None:
            run.font.color.rgb = accent_color
            run.font.bold = True
        else:
            run.font.color.rgb = font_color
            run.font.bold = bold or seg_bold
        run.font.italic = seg_italic


def _add_text_box(slide, text, left, top, width, height,
                  font_name="Calibri", font_size=None,
                  font_color="#000000", alignment=PP_ALIGN.LEFT,
                  bold=False, parse_markdown=True, accent_color=None):
    """Add a text box shape with styled text and design token margins."""
    if font_size is None:
        font_size = _DESIGN_TOKENS["body_size"]
    scaled_size = _scaled_font_size(font_size, font_name)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    # Apply margins from design tokens
    tf.margin_left = _DESIGN_TOKENS["margin_left"]
    tf.margin_right = _DESIGN_TOKENS["margin_right"]
    tf.margin_top = _DESIGN_TOKENS["margin_top"]
    tf.margin_bottom = _DESIGN_TOKENS["margin_bottom"]

    p = tf.paragraphs[0]
    p.alignment = alignment
    p.line_spacing = _DESIGN_TOKENS["line_spacing"]
    p.space_after = _DESIGN_TOKENS["para_space_after"]

    if parse_markdown:
        _apply_markdown_runs(
            p, text, font_name=font_name, font_size=scaled_size,
            font_color=font_color, bold=bold, accent_color=accent_color,
        )
    else:
        p.text = str(text)
        p.font.name = font_name
        p.font.size = scaled_size
        p.font.bold = bold
        if isinstance(font_color, str):
            font_color = RGBColor.from_string(font_color.lstrip("#"))
        p.font.color.rgb = font_color


def _add_bullet_list(slide, items, left, top, width, height,
                     font_name="Calibri", font_size=None,
                     font_color="#000000", accent_color=None):
    """Add a text box with bullet-point paragraphs."""
    if font_size is None:
        font_size = _DESIGN_TOKENS["body_small_size"]
    scaled_size = _scaled_font_size(font_size, font_name)
    if isinstance(font_color, str):
        resolved_color = RGBColor.from_string(font_color.lstrip("#"))
    else:
        resolved_color = font_color

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    # Apply margins
    tf.margin_left = _DESIGN_TOKENS["margin_left"]
    tf.margin_right = _DESIGN_TOKENS["margin_right"]
    tf.margin_top = _DESIGN_TOKENS["margin_top"]
    tf.margin_bottom = _DESIGN_TOKENS["margin_bottom"]

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = _DESIGN_TOKENS["bullet_space_after"]
        p.line_spacing = _DESIGN_TOKENS["line_spacing"]
        _apply_markdown_runs(
            p, f"\u2022 {item}", font_name=font_name,
            font_size=scaled_size, font_color=resolved_color,
            accent_color=accent_color,
        )


def _find_element(slide_data, element_type):
    """Find first element of a given type in slide data."""
    for el in slide_data.elements:
        if el.type == element_type:
            return el
    return None


def _render_kicker(slide, slide_data, theme):
    """Render a kicker element above the title (small uppercase text in accent color).

    Returns True if a kicker was rendered, False otherwise.
    """
    kicker_el = _find_element(slide_data, "kicker")
    if not kicker_el or not kicker_el.content:
        return False
    kicker_text = str(kicker_el.content).upper()
    _add_text_box(slide, kicker_text,
                  left=TITLE_LEFT, top=Inches(0.25),
                  width=TITLE_WIDTH, height=Inches(0.3),
                  font_name=theme.font_body,
                  font_size=_DESIGN_TOKENS["kicker_size"],
                  font_color=theme.colors.accent,
                  bold=True, parse_markdown=False)
    return True


_TAKEAWAY_HEIGHT = Inches(0.6)
_TAKEAWAY_GAP = Inches(0.1)  # Gap between body content and takeaway
_TAKEAWAY_RESERVE = _TAKEAWAY_HEIGHT + _TAKEAWAY_GAP  # Total space to reserve


def _has_takeaway(slide_data):
    """Check if slide data contains a non-empty takeaway element."""
    el = _find_element(slide_data, "takeaway")
    return el is not None and bool(el.content)


def _render_takeaway(slide, slide_data, theme, top=None):
    """Render a takeaway element at the bottom of the slide.

    Returns True if a takeaway was rendered, False otherwise.
    """
    takeaway_el = _find_element(slide_data, "takeaway")
    if not takeaway_el or not takeaway_el.content:
        return False
    takeaway_top = top or Inches(6.4)
    _add_text_box(slide, str(takeaway_el.content),
                  left=TITLE_LEFT, top=takeaway_top,
                  width=TITLE_WIDTH, height=_TAKEAWAY_HEIGHT,
                  font_name=theme.font_body,
                  font_size=_DESIGN_TOKENS["takeaway_size"],
                  font_color=theme.colors.secondary,
                  bold=True)
    return True


def _set_slide_background(slide, theme):
    """Set slide background — supports solid and gradient fills."""
    bg_hex = theme.colors.background.lstrip("#")
    if getattr(theme, "background_style", None) == "gradient":
        try:
            primary_hex = theme.colors.primary.lstrip("#")
            fill = slide.background.fill
            fill.gradient()
            fill.gradient_stops[0].color.rgb = RGBColor.from_string(bg_hex)
            fill.gradient_stops[1].color.rgb = RGBColor.from_string(primary_hex)
            fill.gradient_angle = 270
            return
        except Exception:
            pass  # Fall through to solid fill
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor.from_string(bg_hex)


# === Slide Chrome ===

def _add_slide_chrome(slide, theme, slide_number, total_slides):
    """Add progress bar and slide number to a slide."""
    bar_height = Inches(0.04)
    bar_top = _DESIGN_TOKENS["accent_bar_top"]

    # Background bar — full width, light tint
    bg_tint = _blend_color(theme.colors.text_light, theme.colors.background, 0.20)
    bg_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), bar_top, SLIDE_WIDTH, bar_height,
    )
    bg_bar.fill.solid()
    bg_bar.fill.fore_color.rgb = RGBColor.from_string(bg_tint.lstrip("#"))
    bg_bar.line.fill.background()

    # Fill bar — proportional width, accent color
    progress = slide_number / total_slides
    fill_width = int(SLIDE_WIDTH * progress)
    fill_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), bar_top, fill_width, bar_height,
    )
    fill_bar.fill.solid()
    fill_bar.fill.fore_color.rgb = RGBColor.from_string(
        theme.colors.accent.lstrip("#")
    )
    fill_bar.line.fill.background()

    # Slide number — right-aligned (just the number)
    num_box = slide.shapes.add_textbox(
        SLIDE_WIDTH - MARGIN_RIGHT - Inches(1.2), _DESIGN_TOKENS["footer_top"],
        Inches(1.2), _DESIGN_TOKENS["footer_height"],
    )
    tf = num_box.text_frame
    p = tf.paragraphs[0]
    p.text = str(slide_number)
    p.alignment = PP_ALIGN.RIGHT
    p.font.size = _DESIGN_TOKENS["footer_size"]
    p.font.color.rgb = RGBColor.from_string(
        theme.colors.text_light.lstrip("#")
    )


# === Slide Motif ===

def _add_slide_motif(slide, theme):
    """Draw a consistent decorative motif on every slide, determined by theme."""
    motif_type = int(hashlib.md5(theme.id.encode()).hexdigest()[:8], 16) % 3
    primary_tint = _blend_color(theme.colors.primary, theme.colors.background, 0.40)

    if motif_type == 0:
        # Sidebar: 0.35" vertical bar, full height, left edge
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0), Inches(0.35), SLIDE_HEIGHT,
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = RGBColor.from_string(primary_tint.lstrip("#"))
        bar.line.fill.background()
    elif motif_type == 1:
        # Top band: full width, 0.5" tall
        band = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.5),
        )
        band.fill.solid()
        band.fill.fore_color.rgb = RGBColor.from_string(
            theme.colors.primary.lstrip("#")
        )
        band.line.fill.background()
    else:
        # Corner accent: L-shaped pair in bottom-right
        accent_hex = theme.colors.accent.lstrip("#")
        h_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            SLIDE_WIDTH - Inches(2.0), SLIDE_HEIGHT - Inches(0.15),
            Inches(2.0), Inches(0.15),
        )
        h_bar.fill.solid()
        h_bar.fill.fore_color.rgb = RGBColor.from_string(accent_hex)
        h_bar.line.fill.background()
        v_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            SLIDE_WIDTH - Inches(0.15), SLIDE_HEIGHT - Inches(2.0),
            Inches(0.15), Inches(2.0),
        )
        v_bar.fill.solid()
        v_bar.fill.fore_color.rgb = RGBColor.from_string(accent_hex)
        v_bar.line.fill.background()


# === Icon Bullets ===

# Keyword-to-shape mapping for contextual icon bullets
_ICON_SHAPES = {
    "growth": MSO_SHAPE.UP_ARROW,
    "increase": MSO_SHAPE.UP_ARROW,
    "rise": MSO_SHAPE.UP_ARROW,
    "revenue": MSO_SHAPE.UP_ARROW,
    "security": MSO_SHAPE.DIAMOND,
    "protect": MSO_SHAPE.DIAMOND,
    "safe": MSO_SHAPE.DIAMOND,
    "innovation": MSO_SHAPE.LIGHTNING_BOLT,
    "new": MSO_SHAPE.LIGHTNING_BOLT,
    "launch": MSO_SHAPE.LIGHTNING_BOLT,
    "speed": MSO_SHAPE.LIGHTNING_BOLT,
    "fast": MSO_SHAPE.LIGHTNING_BOLT,
    "team": MSO_SHAPE.STAR_5_POINT,
    "people": MSO_SHAPE.STAR_5_POINT,
    "user": MSO_SHAPE.STAR_5_POINT,
    "customer": MSO_SHAPE.STAR_5_POINT,
    "target": MSO_SHAPE.CROSS,
    "goal": MSO_SHAPE.CROSS,
    "focus": MSO_SHAPE.CROSS,
    "data": MSO_SHAPE.HEXAGON,
    "analytics": MSO_SHAPE.HEXAGON,
    "process": MSO_SHAPE.CHEVRON,
    "step": MSO_SHAPE.CHEVRON,
    "phase": MSO_SHAPE.CHEVRON,
    "strategy": MSO_SHAPE.PENTAGON,
    "plan": MSO_SHAPE.PENTAGON,
    "reduce": MSO_SHAPE.DOWN_ARROW,
    "decrease": MSO_SHAPE.DOWN_ARROW,
    "cost": MSO_SHAPE.DOWN_ARROW,
}

# Fallback shape cycle when no keyword matches
_FALLBACK_SHAPES = [
    MSO_SHAPE.OVAL,
    MSO_SHAPE.DIAMOND,
    MSO_SHAPE.HEXAGON,
    MSO_SHAPE.ROUNDED_RECTANGLE,
]


def _pick_icon_shape(text, index=0):
    """Pick an MSO_SHAPE based on keyword scanning of bullet text."""
    text_lower = text.lower()
    for keyword, shape in _ICON_SHAPES.items():
        if keyword in text_lower:
            return shape
    return _FALLBACK_SHAPES[index % len(_FALLBACK_SHAPES)]


def _add_icon_bullet_list(slide, items, left, top, width, height,
                          font_name="Calibri", font_size=None,
                          font_color="#000000", icon_color="#4472C4",
                          accent_color=None):
    """Add a bullet list with small shape icons next to each item."""
    if font_size is None:
        font_size = _DESIGN_TOKENS["body_size"]
    scaled_size = _scaled_font_size(font_size, font_name)
    if isinstance(font_color, str):
        resolved_color = RGBColor.from_string(font_color.lstrip("#"))
    else:
        resolved_color = font_color
    if isinstance(icon_color, str):
        resolved_icon_color = RGBColor.from_string(icon_color.lstrip("#"))
    else:
        resolved_icon_color = icon_color

    icon_size = Inches(0.22)
    icon_left_offset = left
    text_left_offset = left + Inches(0.35)
    text_width = width - Inches(0.35)

    # Dynamic spacing: scale with font size (pt * 0.023 inches gives ~0.41" at 18pt)
    item_spacing_in = font_size.pt * 0.023
    item_spacing_in = max(item_spacing_in, 0.32)  # Floor to prevent cramming

    txBox = slide.shapes.add_textbox(text_left_offset, top, text_width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left = _DESIGN_TOKENS["margin_left"]
    tf.margin_right = _DESIGN_TOKENS["margin_right"]
    tf.margin_top = _DESIGN_TOKENS["margin_top"]
    tf.margin_bottom = _DESIGN_TOKENS["margin_bottom"]

    for i, item in enumerate(items):
        # Add icon shape only if it fits within the body area
        icon_top = top + Inches(i * item_spacing_in + 0.05)
        if icon_top + icon_size <= top + height:
            shape = _pick_icon_shape(str(item), i)
            icon = slide.shapes.add_shape(
                shape, icon_left_offset, icon_top, icon_size, icon_size,
            )
            icon.fill.solid()
            icon.fill.fore_color.rgb = resolved_icon_color
            icon.line.fill.background()

        # Always add text paragraph (let PowerPoint handle overflow/clipping)
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = _DESIGN_TOKENS["bullet_space_after"]
        p.line_spacing = _DESIGN_TOKENS["line_spacing"]
        _apply_markdown_runs(
            p, str(item), font_name=font_name,
            font_size=scaled_size, font_color=resolved_color,
            accent_color=accent_color,
        )


# === Decorative Accents ===

def _add_decorative_accent(slide, theme, slide_index):
    """Add a small decorative visual accent to text-heavy slides.

    Rotates through 3 styles based on slide_index to add visual variety.
    """
    accent_hex = theme.colors.accent.lstrip("#")
    style = slide_index % 3

    if style == 0:
        # Accent dot cluster — 3 small circles in top-right corner
        for i in range(3):
            dot = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                SLIDE_WIDTH - Inches(0.9 + i * 0.25), Inches(0.3),
                Inches(0.15), Inches(0.15),
            )
            dot.fill.solid()
            dot.fill.fore_color.rgb = RGBColor.from_string(accent_hex)
            dot.line.fill.background()
    elif style == 1:
        # Corner bracket — thin L-shape lines in bottom-left
        h_line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.5), SLIDE_HEIGHT - Inches(0.7),
            Inches(1.2), Inches(0.04),
        )
        h_line.fill.solid()
        h_line.fill.fore_color.rgb = RGBColor.from_string(accent_hex)
        h_line.line.fill.background()
        v_line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.5), SLIDE_HEIGHT - Inches(1.7),
            Inches(0.04), Inches(1.04),
        )
        v_line.fill.solid()
        v_line.fill.fore_color.rgb = RGBColor.from_string(accent_hex)
        v_line.line.fill.background()
    else:
        # Horizontal divider rule below title
        rule = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            BODY_LEFT, Inches(1.65), BODY_WIDTH, Inches(0.03),
        )
        rule.fill.solid()
        rule.fill.fore_color.rgb = RGBColor.from_string(accent_hex)
        rule.line.fill.background()


# === Card Layouts ===

def _add_card(slide, *, left, top, width, height, theme, color_key="primary"):
    """Add a rectangle card with subtle border and left accent strip."""
    # Card background — 10% tint of theme color
    base_color = getattr(theme.colors, color_key)
    card_fill = _blend_color(base_color, theme.colors.background, 0.10)
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height,
    )
    card.adjustments[0] = 0.05  # Corner radius
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor.from_string(card_fill.lstrip("#"))
    card.line.fill.background()  # No visible border

    # Soft drop shadow via XML
    from pptx.oxml.ns import qn
    sp_pr = card._element.spPr
    effect_lst = sp_pr.makeelement(qn("a:effectLst"), {})
    outer_shdw = effect_lst.makeelement(qn("a:outerShdw"), {
        "blurRad": "190500",   # ~15pt blur
        "dist": "63500",       # ~5pt distance
        "dir": "5400000",      # 90 degrees (straight down)
        "algn": "bl",
    })
    srgb = outer_shdw.makeelement(qn("a:srgbClr"), {"val": "000000"})
    alpha = srgb.makeelement(qn("a:alpha"), {"val": "15000"})  # 15% opacity
    srgb.append(alpha)
    outer_shdw.append(srgb)
    effect_lst.append(outer_shdw)
    sp_pr.append(effect_lst)

    # Left accent strip
    strip = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        left, top, Inches(0.06), height,
    )
    strip.fill.solid()
    strip.fill.fore_color.rgb = RGBColor.from_string(
        theme.colors.accent.lstrip("#")
    )
    strip.line.fill.background()

    return card


# === Chart Rendering ===

def _build_chart_palette(theme):
    """Build a 6-color chart palette from theme colors."""
    colors = [theme.colors.accent, theme.colors.primary, theme.colors.secondary]
    # Add 40% tints toward white
    for c in [theme.colors.accent, theme.colors.primary, theme.colors.secondary]:
        colors.append(_blend_color(c, "#FFFFFF", 0.60))
    return colors


def _style_chart(chart, spec, theme):
    """Apply theme styling to a python-pptx chart object."""
    try:
        chart.chart_style = 2  # Minimal style
        plot = chart.plots[0]
        plot.format.fill.background()  # Transparent plot area
    except Exception:
        pass

    # Series colors
    palette = _build_chart_palette(theme)
    is_line_chart = spec.chart_type.value in ("line",)
    try:
        for idx, series in enumerate(chart.series):
            color_hex = palette[idx % len(palette)].lstrip("#")
            series.format.fill.solid()
            series.format.fill.fore_color.rgb = RGBColor.from_string(color_hex)
            # Thicker lines for line charts
            if is_line_chart:
                series.format.line.width = Pt(3)
                series.format.line.color.rgb = RGBColor.from_string(color_hex)
    except Exception:
        pass

    # Data labels on first series for line charts
    if is_line_chart:
        try:
            chart.series[0].has_data_labels = True
            dl = chart.series[0].data_labels
            dl.font.size = _DESIGN_TOKENS["code_size"]
            dl.font.color.rgb = RGBColor.from_string(
                theme.colors.text_light.lstrip("#")
            )
        except Exception:
            pass

    # Axis labels
    try:
        if hasattr(chart, 'value_axis') and chart.value_axis is not None:
            va = chart.value_axis
            va.has_major_gridlines = True
            va.major_gridlines.format.line.color.rgb = RGBColor.from_string(
                theme.colors.text_light.lstrip("#")
            )
            va.format.line.fill.background()  # Hide redundant axis line
            if va.has_tick_labels:
                va.tick_labels.font.size = _DESIGN_TOKENS["body_small_size"]
                va.tick_labels.font.color.rgb = RGBColor.from_string(
                    theme.colors.text_light.lstrip("#")
                )
    except Exception:
        pass

    try:
        if hasattr(chart, 'category_axis') and chart.category_axis is not None:
            ca = chart.category_axis
            ca.has_major_gridlines = False
            if ca.has_tick_labels:
                ca.tick_labels.font.size = _DESIGN_TOKENS["body_small_size"]
                ca.tick_labels.font.color.rgb = RGBColor.from_string(
                    theme.colors.text_light.lstrip("#")
                )
    except Exception:
        pass

    # Legend
    try:
        if chart.has_legend and chart.legend is not None:
            chart.legend.font.size = _DESIGN_TOKENS["code_size"]
    except Exception:
        pass


def _add_chart(slide, spec, theme):
    """Add a native python-pptx chart to the slide. Returns True on success."""
    if not _CHARTS_AVAILABLE:
        return False
    try:
        from core.schemas.studio_schema import ChartType

        # Chart container card — rounded rect behind chart area
        card_margin = Inches(0.15)
        _add_card(slide,
                  left=CHART_LEFT - card_margin,
                  top=CHART_TOP - card_margin,
                  width=CHART_WIDTH + card_margin * 2,
                  height=CHART_HEIGHT + card_margin * 2,
                  theme=theme, color_key="primary")

        if spec.chart_type.value == "scatter":
            chart_data = XyChartData()
            series = chart_data.add_series(spec.title or "Data")
            for pt in spec.points:
                series.add_data_point(pt.x, pt.y)
            chart_type = XL_CHART_TYPE.XY_SCATTER
        else:
            chart_data = CategoryChartData()
            chart_data.categories = spec.categories
            series_list = spec.series
            if spec.chart_type == ChartType.pie and len(series_list) > 1:
                series_list = series_list[:1]
            for s in series_list:
                chart_data.add_series(s.name, s.values)
            chart_type = _CHART_TYPE_MAP.get(spec.chart_type.value, XL_CHART_TYPE.COLUMN_CLUSTERED)

        chart_frame = slide.shapes.add_chart(
            chart_type, CHART_LEFT, CHART_TOP, CHART_WIDTH, CHART_HEIGHT, chart_data
        )
        chart = chart_frame.chart
        chart.has_legend = len(spec.series) > 1 or spec.chart_type.value == "pie"

        _style_chart(chart, spec, theme)

        # Enable data labels on bar/pie charts with <=6 categories
        cat_count = len(spec.categories) if spec.categories else 0
        if spec.chart_type.value in ("bar", "pie") and cat_count <= 6:
            try:
                for series_obj in chart.series:
                    series_obj.has_data_labels = True
                    dl = series_obj.data_labels
                    dl.font.size = _DESIGN_TOKENS["code_size"]
                    dl.font.color.rgb = RGBColor.from_string(
                        theme.colors.text.lstrip("#")
                    )
            except Exception:
                pass

        return True
    except Exception:
        return False


# === Slide-Type Renderer Functions ===

def _render_title(slide, slide_data, theme, **kwargs):
    """Title slide: centered title + subtitle, with optional dark background."""
    title_bg = getattr(theme.colors, "title_background", None)
    if title_bg:
        # Dark title slide — override background and use white text
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor.from_string(title_bg.lstrip("#"))
        title_color = "#FFFFFF"
        subtitle_color = "#CCCCCC"
    else:
        _set_slide_background(slide, theme)
        title_color = theme.colors.primary
        subtitle_color = theme.colors.text_light

    # Decorative banner behind title
    banner_fill = _blend_color(theme.colors.accent, title_bg or theme.colors.background, 0.15)
    banner = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1.5), Inches(2.0), Inches(10.333), Inches(2.5),
    )
    banner.fill.solid()
    banner.fill.fore_color.rgb = RGBColor.from_string(banner_fill.lstrip("#"))
    banner.line.fill.background()

    _add_text_box(slide, slide_data.title or "",
                  left=MARGIN_LEFT, top=Inches(2.5),
                  width=CONTENT_WIDTH, height=Inches(1.5),
                  font_name=theme.font_heading, font_size=_DESIGN_TOKENS["title_size"],
                  font_color=title_color, alignment=PP_ALIGN.CENTER,
                  bold=True)
    subtitle_el = _find_element(slide_data, "subtitle")
    if subtitle_el and subtitle_el.content:
        _add_text_box(slide, subtitle_el.content,
                      left=MARGIN_LEFT, top=Inches(4.2),
                      width=CONTENT_WIDTH, height=Inches(0.8),
                      font_name=theme.font_body, font_size=_DESIGN_TOKENS["subheading_size"],
                      font_color=subtitle_color, alignment=PP_ALIGN.CENTER)


def _render_content(slide, slide_data, theme, **kwargs):
    """Standard content slide: title + body/bullets with varied card styles."""
    slide_index = kwargs.get("slide_index", 0)
    card_style = slide_index % 3

    _render_kicker(slide, slide_data, theme)

    body_top, title_font = _compute_body_top(slide_data.title)
    body_height = SLIDE_HEIGHT - body_top - MARGIN_BOTTOM
    if _has_takeaway(slide_data):
        body_height -= _TAKEAWAY_RESERVE

    _add_text_box(slide, slide_data.title or "",
                  left=TITLE_LEFT, top=TITLE_TOP,
                  width=TITLE_WIDTH, height=TITLE_HEIGHT,
                  font_name=theme.font_heading, font_size=title_font,
                  font_color=theme.colors.primary, bold=True)

    body_el = _find_element(slide_data, "body")
    bullet_el = _find_element(slide_data, "bullet_list")

    content_left = BODY_LEFT
    content_width = BODY_WIDTH

    if card_style == 0:
        # Style 0: Full card with left accent strip (original)
        _add_card(slide, left=BODY_LEFT, top=body_top,
                  width=BODY_WIDTH, height=body_height,
                  theme=theme, color_key="primary")
        content_left = BODY_LEFT + Inches(0.12)
        content_width = BODY_WIDTH - Inches(0.12)
    elif card_style == 1:
        # Style 1: Top accent bar above content area
        top_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            BODY_LEFT, body_top, BODY_WIDTH, Inches(0.08),
        )
        top_bar.fill.solid()
        top_bar.fill.fore_color.rgb = RGBColor.from_string(
            theme.colors.accent.lstrip("#")
        )
        top_bar.line.fill.background()
    else:
        # Style 2: Horizontal rule below title
        rule_top = body_top - Inches(0.15)
        rule = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            BODY_LEFT, rule_top, BODY_WIDTH, Inches(0.03),
        )
        rule.fill.solid()
        rule.fill.fore_color.rgb = RGBColor.from_string(
            theme.colors.accent.lstrip("#")
        )
        rule.line.fill.background()

    if bullet_el and isinstance(bullet_el.content, list):
        _add_icon_bullet_list(slide, bullet_el.content,
                              left=content_left, top=body_top,
                              width=content_width, height=body_height,
                              font_name=theme.font_body, font_size=_DESIGN_TOKENS["body_size"],
                              font_color=theme.colors.text,
                              icon_color=theme.colors.accent,
                              accent_color=theme.colors.accent)
    elif body_el and body_el.content:
        _add_text_box(slide, body_el.content,
                      left=content_left, top=body_top,
                      width=content_width, height=body_height,
                      font_name=theme.font_body, font_size=_DESIGN_TOKENS["body_size"],
                      font_color=theme.colors.text,
                      accent_color=theme.colors.accent)

    _render_takeaway(slide, slide_data, theme)
    _set_slide_background(slide, theme)


def _render_two_column(slide, slide_data, theme, **kwargs):
    """Two-column layout: title + left/right body areas."""
    _render_kicker(slide, slide_data, theme)

    body_top, title_font = _compute_body_top(slide_data.title)
    body_height = SLIDE_HEIGHT - body_top - MARGIN_BOTTOM
    if _has_takeaway(slide_data):
        body_height -= _TAKEAWAY_RESERVE

    _add_text_box(slide, slide_data.title or "",
                  left=TITLE_LEFT, top=TITLE_TOP,
                  width=TITLE_WIDTH, height=TITLE_HEIGHT,
                  font_name=theme.font_heading, font_size=title_font,
                  font_color=theme.colors.primary, bold=True)

    body_elements = [el for el in slide_data.elements if el.type == "body"]
    bullet_elements = [el for el in slide_data.elements if el.type == "bullet_list"]

    # Left column
    left_content = body_elements[0].content if body_elements else ""
    if bullet_elements and isinstance(bullet_elements[0].content, list):
        _add_icon_bullet_list(slide, bullet_elements[0].content,
                              left=MARGIN_LEFT, top=body_top,
                              width=COLUMN_WIDTH, height=body_height,
                              font_name=theme.font_body, font_size=_DESIGN_TOKENS["body_small_size"],
                              font_color=theme.colors.text,
                              icon_color=theme.colors.accent,
                              accent_color=theme.colors.accent)
    else:
        _add_text_box(slide, left_content,
                      left=MARGIN_LEFT, top=body_top,
                      width=COLUMN_WIDTH, height=body_height,
                      font_name=theme.font_body, font_size=_DESIGN_TOKENS["body_small_size"],
                      font_color=theme.colors.text,
                      accent_color=theme.colors.accent)

    # Right column
    right_content = body_elements[1].content if len(body_elements) > 1 else ""
    if len(bullet_elements) > 1 and isinstance(bullet_elements[1].content, list):
        _add_icon_bullet_list(slide, bullet_elements[1].content,
                              left=MARGIN_LEFT + COLUMN_WIDTH + COLUMN_GAP, top=body_top,
                              width=COLUMN_WIDTH, height=body_height,
                              font_name=theme.font_body, font_size=_DESIGN_TOKENS["body_small_size"],
                              font_color=theme.colors.text,
                              icon_color=theme.colors.accent,
                              accent_color=theme.colors.accent)
    else:
        _add_text_box(slide, right_content,
                      left=MARGIN_LEFT + COLUMN_WIDTH + COLUMN_GAP, top=body_top,
                      width=COLUMN_WIDTH, height=body_height,
                      font_name=theme.font_body, font_size=_DESIGN_TOKENS["body_small_size"],
                      font_color=theme.colors.text,
                      accent_color=theme.colors.accent)

    _render_takeaway(slide, slide_data, theme)
    _add_decorative_accent(slide, theme, kwargs.get("slide_index", 0))
    _set_slide_background(slide, theme)


def _render_comparison(slide, slide_data, theme, **kwargs):
    """Comparison slide: title + two labeled columns in cards."""
    _render_kicker(slide, slide_data, theme)

    body_top, title_font = _compute_body_top(slide_data.title)
    body_height = SLIDE_HEIGHT - body_top - MARGIN_BOTTOM
    if _has_takeaway(slide_data):
        body_height -= _TAKEAWAY_RESERVE

    _add_text_box(slide, slide_data.title or "",
                  left=TITLE_LEFT, top=TITLE_TOP,
                  width=TITLE_WIDTH, height=TITLE_HEIGHT,
                  font_name=theme.font_heading, font_size=title_font,
                  font_color=theme.colors.primary, bold=True)

    body_elements = [el for el in slide_data.elements if el.type == "body"]
    left_text = body_elements[0].content if body_elements else ""
    right_text = body_elements[1].content if len(body_elements) > 1 else ""

    # Left card (primary)
    _add_card(slide, left=MARGIN_LEFT, top=body_top,
              width=COLUMN_WIDTH, height=body_height,
              theme=theme, color_key="primary")
    _add_text_box(slide, left_text,
                  left=MARGIN_LEFT + Inches(0.12), top=body_top,
                  width=COLUMN_WIDTH - Inches(0.12), height=body_height,
                  font_name=theme.font_body, font_size=_DESIGN_TOKENS["body_small_size"],
                  font_color=theme.colors.text,
                  accent_color=theme.colors.accent)

    # Right card (secondary)
    right_left = MARGIN_LEFT + COLUMN_WIDTH + COLUMN_GAP
    _add_card(slide, left=right_left, top=body_top,
              width=COLUMN_WIDTH, height=body_height,
              theme=theme, color_key="secondary")
    _add_text_box(slide, right_text,
                  left=right_left + Inches(0.12), top=body_top,
                  width=COLUMN_WIDTH - Inches(0.12), height=body_height,
                  font_name=theme.font_body, font_size=_DESIGN_TOKENS["body_small_size"],
                  font_color=theme.colors.text,
                  accent_color=theme.colors.accent)

    _render_takeaway(slide, slide_data, theme)
    _add_decorative_accent(slide, theme, kwargs.get("slide_index", 0))
    _set_slide_background(slide, theme)


def _render_timeline(slide, slide_data, theme, **kwargs):
    """Timeline/roadmap slide: title + sequential items with visual circles."""
    _render_kicker(slide, slide_data, theme)

    body_top, title_font = _compute_body_top(slide_data.title)
    body_height = SLIDE_HEIGHT - body_top - MARGIN_BOTTOM
    if _has_takeaway(slide_data):
        body_height -= _TAKEAWAY_RESERVE

    _add_text_box(slide, slide_data.title or "",
                  left=TITLE_LEFT, top=TITLE_TOP,
                  width=TITLE_WIDTH, height=TITLE_HEIGHT,
                  font_name=theme.font_heading, font_size=title_font,
                  font_color=theme.colors.primary, bold=True)

    bullet_el = _find_element(slide_data, "bullet_list")
    body_el = _find_element(slide_data, "body")

    if bullet_el and isinstance(bullet_el.content, list):
        # Draw accent circles for each timeline item
        item_count = len(bullet_el.content)
        for idx in range(item_count):
            circle_top = body_top + Inches(idx * 0.55 + 0.1)
            circle = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                BODY_LEFT, circle_top, Inches(0.25), Inches(0.25),
            )
            circle.fill.solid()
            circle.fill.fore_color.rgb = RGBColor.from_string(
                theme.colors.accent.lstrip("#")
            )
            circle.line.fill.background()

        # Shift bullet text right to accommodate circles
        _add_bullet_list(slide, bullet_el.content,
                         left=BODY_LEFT + Inches(0.4), top=body_top,
                         width=BODY_WIDTH - Inches(0.4), height=body_height,
                         font_name=theme.font_body, font_size=_DESIGN_TOKENS["body_size"],
                         font_color=theme.colors.text,
                         accent_color=theme.colors.accent)
    elif body_el and body_el.content:
        _add_text_box(slide, body_el.content,
                      left=BODY_LEFT, top=body_top,
                      width=BODY_WIDTH, height=body_height,
                      font_name=theme.font_body, font_size=_DESIGN_TOKENS["body_size"],
                      font_color=theme.colors.text,
                      accent_color=theme.colors.accent)

    _render_takeaway(slide, slide_data, theme)
    _set_slide_background(slide, theme)


def _render_chart(slide, slide_data, theme, **kwargs):
    """Chart slide: title + native chart or text fallback."""
    from core.studio.slides.charts import parse_chart_spec, normalize_chart_spec

    _render_kicker(slide, slide_data, theme)

    _, title_font = _compute_body_top(slide_data.title)

    _add_text_box(slide, slide_data.title or "",
                  left=TITLE_LEFT, top=TITLE_TOP,
                  width=TITLE_WIDTH, height=TITLE_HEIGHT,
                  font_name=theme.font_heading, font_size=title_font,
                  font_color=theme.colors.primary, bold=True)

    chart_el = _find_element(slide_data, "chart")
    body_el = _find_element(slide_data, "body")
    chart_rendered = False

    if chart_el and chart_el.content:
        spec = parse_chart_spec(chart_el.content)
        if spec is not None:
            spec = normalize_chart_spec(spec)
            chart_rendered = _add_chart(slide, spec, theme)

    # Fallback: text placeholder if chart parse/render failed
    if not chart_rendered and chart_el and chart_el.content:
        content = chart_el.content if isinstance(chart_el.content, str) else str(chart_el.content)
        _add_text_box(slide, f"[Chart: {content}]",
                      left=BODY_LEFT, top=BODY_TOP,
                      width=BODY_WIDTH, height=Inches(3.0),
                      font_name=theme.font_body, font_size=_DESIGN_TOKENS["body_small_size"],
                      font_color=theme.colors.secondary,
                      alignment=PP_ALIGN.CENTER)

    # Body/caption below chart
    if body_el and body_el.content:
        caption_top = CAPTION_TOP if chart_rendered else Inches(5.0)
        caption_h = CAPTION_HEIGHT
        if _has_takeaway(slide_data):
            caption_h = Inches(0.3)  # Shrink to leave room for takeaway
        _add_text_box(slide, body_el.content,
                      left=BODY_LEFT, top=caption_top,
                      width=BODY_WIDTH, height=caption_h,
                      font_name=theme.font_body, font_size=_DESIGN_TOKENS["body_small_size"],
                      font_color=theme.colors.text)

    _render_takeaway(slide, slide_data, theme)
    _set_slide_background(slide, theme)


def _render_image_text(slide, slide_data, theme, **kwargs):
    """Image+text slide: split layout with image (or placeholder) and body."""
    body_top, title_font = _compute_body_top(slide_data.title)
    body_height = SLIDE_HEIGHT - body_top - MARGIN_BOTTOM

    _add_text_box(slide, slide_data.title or "",
                  left=TITLE_LEFT, top=TITLE_TOP,
                  width=TITLE_WIDTH, height=TITLE_HEIGHT,
                  font_name=theme.font_heading, font_size=title_font,
                  font_color=theme.colors.primary, bold=True)

    image_el = _find_element(slide_data, "image")
    body_el = _find_element(slide_data, "body")

    # Check for a generated image
    images = kwargs.get("images") or {}
    img_buf = images.get(slide_data.id) if slide_data.id else None

    if img_buf is not None:
        # Embed actual image
        img_buf.seek(0)
        pic = slide.shapes.add_picture(
            img_buf, MARGIN_LEFT, body_top, COLUMN_WIDTH, body_height,
        )
        # Add 2pt border in primary color
        pic.line.width = Pt(2)
        pic.line.color.rgb = RGBColor.from_string(
            theme.colors.primary.lstrip("#")
        )
    else:
        # Text placeholder fallback
        placeholder_text = "[Image]"
        if image_el and image_el.content:
            placeholder_text = f"[Image: {image_el.content}]"
        _add_text_box(slide, placeholder_text,
                      left=MARGIN_LEFT, top=body_top,
                      width=COLUMN_WIDTH, height=body_height,
                      font_name=theme.font_body, font_size=_DESIGN_TOKENS["body_small_size"],
                      font_color=theme.colors.secondary,
                      alignment=PP_ALIGN.CENTER)

    if body_el and body_el.content:
        _add_text_box(slide, body_el.content,
                      left=MARGIN_LEFT + COLUMN_WIDTH + COLUMN_GAP, top=body_top,
                      width=COLUMN_WIDTH, height=body_height,
                      font_name=theme.font_body, font_size=_DESIGN_TOKENS["body_small_size"],
                      font_color=theme.colors.text,
                      accent_color=theme.colors.accent)

    _set_slide_background(slide, theme)


def _render_quote(slide, slide_data, theme, **kwargs):
    """Quote slide: large quote text with accent bar and attribution."""
    quote_el = _find_element(slide_data, "quote")
    body_el = _find_element(slide_data, "body")

    quote_text = ""
    if quote_el and quote_el.content:
        raw = quote_el.content.strip().strip('"\u201C\u201D')
        quote_text = f"\u201C{raw}\u201D"
    elif slide_data.title:
        quote_text = slide_data.title

    # Vertical accent bar to the left of quote
    accent_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(1.2), Inches(1.5), Inches(0.12), Inches(4.0),
    )
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = RGBColor.from_string(
        theme.colors.accent.lstrip("#")
    )
    accent_bar.line.fill.background()

    _add_text_box(slide, quote_text,
                  left=Inches(1.5), top=Inches(1.5),
                  width=Inches(10.333), height=Inches(4.0),
                  font_name=theme.font_heading, font_size=_DESIGN_TOKENS["quote_size"],
                  font_color=theme.colors.primary,
                  alignment=PP_ALIGN.CENTER)

    if body_el and body_el.content:
        _add_text_box(slide, f"\u2014 {body_el.content}",
                      left=Inches(1.5), top=Inches(5.8),
                      width=Inches(10.333), height=Inches(0.8),
                      font_name=theme.font_body, font_size=_DESIGN_TOKENS["attribution_size"],
                      font_color=theme.colors.text_light,
                      alignment=PP_ALIGN.CENTER)

    _set_slide_background(slide, theme)


def _render_code(slide, slide_data, theme, **kwargs):
    """Code slide: title + monospace code block."""
    _add_text_box(slide, slide_data.title or "",
                  left=TITLE_LEFT, top=TITLE_TOP,
                  width=TITLE_WIDTH, height=TITLE_HEIGHT,
                  font_name=theme.font_heading, font_size=_DESIGN_TOKENS["heading_size"],
                  font_color=theme.colors.primary, bold=True)

    code_el = _find_element(slide_data, "code")
    if code_el and code_el.content:
        _add_text_box(slide, code_el.content,
                      left=Inches(1.0), top=BODY_TOP,
                      width=Inches(11.333), height=BODY_HEIGHT,
                      font_name="Courier New", font_size=_DESIGN_TOKENS["code_size"],
                      font_color=theme.colors.text, parse_markdown=False)

    _set_slide_background(slide, theme)


def _render_team(slide, slide_data, theme, **kwargs):
    """Team/credits slide: title + team member list."""
    body_top, title_font = _compute_body_top(slide_data.title)
    body_height = SLIDE_HEIGHT - body_top - MARGIN_BOTTOM

    _add_text_box(slide, slide_data.title or "",
                  left=TITLE_LEFT, top=TITLE_TOP,
                  width=TITLE_WIDTH, height=TITLE_HEIGHT,
                  font_name=theme.font_heading, font_size=title_font,
                  font_color=theme.colors.primary, bold=True)

    bullet_el = _find_element(slide_data, "bullet_list")
    body_el = _find_element(slide_data, "body")

    if bullet_el and isinstance(bullet_el.content, list):
        _add_bullet_list(slide, bullet_el.content,
                         left=BODY_LEFT, top=body_top,
                         width=BODY_WIDTH, height=body_height,
                         font_name=theme.font_body, font_size=_DESIGN_TOKENS["body_size"],
                         font_color=theme.colors.text,
                         accent_color=theme.colors.accent)
    elif body_el and body_el.content:
        _add_text_box(slide, body_el.content,
                      left=BODY_LEFT, top=body_top,
                      width=BODY_WIDTH, height=body_height,
                      font_name=theme.font_body, font_size=_DESIGN_TOKENS["body_size"],
                      font_color=theme.colors.text,
                      accent_color=theme.colors.accent)

    _add_decorative_accent(slide, theme, kwargs.get("slide_index", 0))
    _set_slide_background(slide, theme)


def _render_stat(slide, slide_data, theme, **kwargs):
    """Stat slide: 1-3 large stat callouts with labels."""
    _render_kicker(slide, slide_data, theme)

    _, title_font = _compute_body_top(slide_data.title)

    _add_text_box(slide, slide_data.title or "",
                  left=TITLE_LEFT, top=TITLE_TOP,
                  width=TITLE_WIDTH, height=TITLE_HEIGHT,
                  font_name=theme.font_heading, font_size=title_font,
                  font_color=theme.colors.primary, bold=True)

    stat_el = _find_element(slide_data, "stat_callout")
    body_el = _find_element(slide_data, "body")

    stats = []
    if stat_el and isinstance(stat_el.content, list):
        for item in stat_el.content[:3]:
            if isinstance(item, dict) and "value" in item:
                stats.append(item)
    elif stat_el and isinstance(stat_el.content, str):
        # Fallback: try parsing JSON string
        try:
            import json
            parsed = json.loads(stat_el.content)
            if isinstance(parsed, list):
                for item in parsed[:3]:
                    if isinstance(item, dict) and "value" in item:
                        stats.append(item)
        except (json.JSONDecodeError, TypeError):
            pass

    if stats:
        col_count = len(stats)
        col_width = CONTENT_WIDTH / col_count
        for idx, stat in enumerate(stats):
            col_left = MARGIN_LEFT + col_width * idx
            # Large value
            _add_text_box(slide, stat.get("value", ""),
                          left=col_left, top=Inches(2.5),
                          width=col_width, height=Inches(2.0),
                          font_name=theme.font_heading,
                          font_size=_DESIGN_TOKENS["stat_callout_size"],
                          font_color=theme.colors.accent,
                          alignment=PP_ALIGN.CENTER, bold=True)
            # Label below
            _add_text_box(slide, stat.get("label", ""),
                          left=col_left, top=Inches(4.5),
                          width=col_width, height=Inches(0.8),
                          font_name=theme.font_body,
                          font_size=_DESIGN_TOKENS["stat_label_size"],
                          font_color=theme.colors.text_light,
                          alignment=PP_ALIGN.CENTER)
    elif body_el and body_el.content:
        # Graceful fallback: render as body text
        _add_text_box(slide, body_el.content,
                      left=BODY_LEFT, top=BODY_TOP,
                      width=BODY_WIDTH, height=BODY_HEIGHT,
                      font_name=theme.font_body, font_size=_DESIGN_TOKENS["body_size"],
                      font_color=theme.colors.text)

    # Optional body/context below stats
    if stats and body_el and body_el.content:
        stat_body_h = Inches(1.2)
        if _has_takeaway(slide_data):
            stat_body_h = Inches(0.5)  # Shrink to leave room for takeaway
        _add_text_box(slide, body_el.content,
                      left=BODY_LEFT, top=Inches(5.5),
                      width=BODY_WIDTH, height=stat_body_h,
                      font_name=theme.font_body, font_size=_DESIGN_TOKENS["body_small_size"],
                      font_color=theme.colors.text_light,
                      alignment=PP_ALIGN.CENTER)

    _render_takeaway(slide, slide_data, theme)
    _set_slide_background(slide, theme)


def _render_image_full(slide, slide_data, theme, **kwargs):
    """Full-bleed image slide: image covers entire slide with dark overlay and white text."""
    images = kwargs.get("images") or {}
    img_buf = images.get(slide_data.id) if slide_data.id else None

    if img_buf is not None:
        # Embed full-bleed image covering entire slide
        img_buf.seek(0)
        slide.shapes.add_picture(
            img_buf, Inches(0), Inches(0), SLIDE_WIDTH, SLIDE_HEIGHT,
        )
    else:
        # Dark background fallback when no image available
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor.from_string("1A1A2E")

    # Gradient dark overlay: transparent at top, dark at bottom (where text sits)
    overlay = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), SLIDE_WIDTH, SLIDE_HEIGHT,
    )
    overlay.fill.gradient()
    overlay.line.fill.background()

    # Configure gradient stops via XML for precise alpha control
    from pptx.oxml.ns import qn
    sp_pr = overlay._element.spPr
    grad_fill = sp_pr.find(qn("a:gradFill"))
    if grad_fill is not None:
        # Set linear angle: 16200000 = 270 degrees (bottom-to-top in EMU)
        lin = grad_fill.find(qn("a:lin"))
        if lin is None:
            lin = grad_fill.makeelement(qn("a:lin"), {})
            grad_fill.append(lin)
        lin.set("ang", "16200000")
        lin.set("scaled", "1")

        # Replace auto-generated stops with our custom ones
        gs_lst = grad_fill.find(qn("a:gsLst"))
        if gs_lst is not None:
            for gs in list(gs_lst):
                gs_lst.remove(gs)

            # Stop 0 (top of slide): black at 0% opacity (fully transparent)
            gs0 = gs_lst.makeelement(qn("a:gs"), {"pos": "0"})
            srgb0 = gs0.makeelement(qn("a:srgbClr"), {"val": "000000"})
            alpha0 = srgb0.makeelement(qn("a:alpha"), {"val": "0"})
            srgb0.append(alpha0)
            gs0.append(srgb0)
            gs_lst.append(gs0)

            # Stop 1 (bottom of slide): black at 80% opacity
            gs1 = gs_lst.makeelement(qn("a:gs"), {"pos": "100000"})
            srgb1 = gs1.makeelement(qn("a:srgbClr"), {"val": "000000"})
            alpha1 = srgb1.makeelement(qn("a:alpha"), {"val": "80000"})
            srgb1.append(alpha1)
            gs1.append(srgb1)
            gs_lst.append(gs1)

    # Title — centered, white, large
    _add_text_box(slide, slide_data.title or "",
                  left=Inches(1.5), top=Inches(2.0),
                  width=Inches(10.333), height=Inches(2.0),
                  font_name=theme.font_heading, font_size=_DESIGN_TOKENS["title_size"],
                  font_color="#FFFFFF", alignment=PP_ALIGN.CENTER,
                  bold=True)

    # Optional body text below title
    body_el = _find_element(slide_data, "body")
    if body_el and body_el.content:
        _add_text_box(slide, body_el.content,
                      left=Inches(2.0), top=Inches(4.5),
                      width=Inches(9.333), height=Inches(2.0),
                      font_name=theme.font_body, font_size=_DESIGN_TOKENS["subheading_size"],
                      font_color="#DDDDDD", alignment=PP_ALIGN.CENTER)


def _render_section_divider(slide, slide_data, theme, **kwargs):
    """Section divider slide: large section number + title, whitespace-heavy."""
    section_number = kwargs.get("section_number", 1)

    # Use dark background like title slide
    title_bg = getattr(theme.colors, "title_background", None)
    if title_bg:
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor.from_string(title_bg.lstrip("#"))
        title_color = "#FFFFFF"
        subtitle_color = "#CCCCCC"
    else:
        _set_slide_background(slide, theme)
        title_color = theme.colors.primary
        subtitle_color = theme.colors.text_light

    # Large section number in accent color
    section_num = str(section_number)
    _add_text_box(slide, section_num,
                  left=MARGIN_LEFT, top=Inches(1.5),
                  width=CONTENT_WIDTH, height=Inches(2.0),
                  font_name=theme.font_heading,
                  font_size=Pt(120),
                  font_color=theme.colors.accent,
                  alignment=PP_ALIGN.CENTER,
                  bold=True, parse_markdown=False)

    # Section title
    _add_text_box(slide, slide_data.title or "",
                  left=MARGIN_LEFT, top=Inches(3.8),
                  width=CONTENT_WIDTH, height=Inches(1.5),
                  font_name=theme.font_heading,
                  font_size=_DESIGN_TOKENS["heading_size"],
                  font_color=title_color,
                  alignment=PP_ALIGN.CENTER,
                  bold=True)

    # Optional subtitle
    subtitle_el = _find_element(slide_data, "subtitle")
    if subtitle_el and subtitle_el.content:
        _add_text_box(slide, subtitle_el.content,
                      left=MARGIN_LEFT, top=Inches(5.3),
                      width=CONTENT_WIDTH, height=Inches(0.8),
                      font_name=theme.font_body,
                      font_size=_DESIGN_TOKENS["body_size"],
                      font_color=subtitle_color,
                      alignment=PP_ALIGN.CENTER)


# Renderer dispatch table
_RENDERERS = {
    "title": _render_title,
    "content": _render_content,
    "two_column": _render_two_column,
    "comparison": _render_comparison,
    "timeline": _render_timeline,
    "chart": _render_chart,
    "image_text": _render_image_text,
    "image_full": _render_image_full,
    "quote": _render_quote,
    "code": _render_code,
    "team": _render_team,
    "stat": _render_stat,
    "section_divider": _render_section_divider,
}
