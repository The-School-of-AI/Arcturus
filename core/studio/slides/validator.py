"""Open-validation for exported PPTX files."""

from pathlib import Path
from typing import Any, Dict


def validate_pptx(file_path: Path, expected_slide_count: int | None = None) -> Dict[str, Any]:
    """Validate a PPTX file by reloading and checking structural properties.

    Returns a dict with valid, slide_count, has_notes, errors,
    layout_valid, layout_warnings.
    """
    errors = []
    slide_count = 0
    has_notes = False

    try:
        from pptx import Presentation
        prs = Presentation(str(file_path))
        slide_count = len(prs.slides)

        if expected_slide_count is not None and slide_count != expected_slide_count:
            errors.append(
                f"Slide count mismatch: expected {expected_slide_count}, got {slide_count}"
            )

        for slide in prs.slides:
            try:
                notes_slide = slide.notes_slide
                if notes_slide and notes_slide.notes_text_frame.text.strip():
                    has_notes = True
                    break
            except Exception:
                continue

    except Exception as e:
        errors.append(f"Failed to open PPTX: {str(e)}")

    # Layout-quality heuristic
    layout_warnings = []
    try:
        from pptx import Presentation as _Prs
        _prs = _Prs(str(file_path))
        BLOCK_CHAR_LIMIT = 800
        SLIDE_CHAR_LIMIT = 1600
        for slide_idx, slide in enumerate(_prs.slides):
            slide_total = 0
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text_len = len(shape.text_frame.text)
                    slide_total += text_len
                    if text_len > BLOCK_CHAR_LIMIT:
                        layout_warnings.append(
                            f"Slide {slide_idx + 1}: text block exceeds {BLOCK_CHAR_LIMIT} chars ({text_len} chars)"
                        )
            if slide_total > SLIDE_CHAR_LIMIT:
                layout_warnings.append(
                    f"Slide {slide_idx + 1}: total text density exceeds {SLIDE_CHAR_LIMIT} chars ({slide_total} chars)"
                )
    except Exception:
        pass

    return {
        "valid": len(errors) == 0,
        "slide_count": slide_count,
        "has_notes": has_notes,
        "errors": errors,
        "layout_valid": len(layout_warnings) == 0,
        "layout_warnings": layout_warnings,
    }
